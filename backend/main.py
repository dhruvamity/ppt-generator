import os
import re
import json
from io import BytesIO
from typing import List, Optional
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from lxml import etree
import latex2mathml.converter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from google import genai
from dotenv import load_dotenv

load_dotenv(os.path.join(os.path.dirname(__file__), '..', '.env'))

app = FastAPI(title="SlideGen Pro Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

THEMES = {
    'dark-neon': {
        'bgColor': '05050F', 'cyan': '00E5FF', 'purple': '7C3AED',
        'gold': 'FFD700', 'tealDecor': '008080', 'bgCard': '0A0A1F', 'textWhite': 'FFFFFF',
        'textBlack': '000000', 'decorPurple': '7C3AED'
    },
    'clean-light': {
        'bgColor': 'F8F9FF', 'cyan': '2563EB', 'purple': '0D9488',
        'gold': 'F59E0B', 'tealDecor': '86F2E4', 'bgCard': 'FFFFFF', 'textWhite': '0B1C30',
        'textBlack': 'FFFFFF', 'decorPurple': '0D9488'
    },
    'minimalist-blue': {
        'bgColor': '0F172A', 'cyan': '38BDF8', 'purple': '818CF8',
        'gold': '94A3B8', 'tealDecor': '1E293B', 'bgCard': '1E293B', 'textWhite': 'F8FAFC',
        'textBlack': '0F172A', 'decorPurple': '818CF8'
    },
    'retro-terminal': {
        'bgColor': '001A00', 'cyan': '00FF00', 'purple': '008800',
        'gold': '00FF44', 'tealDecor': '003300', 'bgCard': '002200', 'textWhite': '00FF00',
        'textBlack': '000000', 'decorPurple': '005500'
    },
    'academic-classic': {
        'bgColor': 'FDFBF7', 'cyan': '800020', 'purple': '002147',
        'gold': 'D4AF37', 'tealDecor': 'E8E3D2', 'bgCard': 'FFFFFF', 'textWhite': '2C2C2C',
        'textBlack': 'FFFFFF', 'decorPurple': '002147'
    },
    'sunset-orange': {
        'bgColor': '1A1A1A', 'cyan': 'FF4500', 'purple': 'FF8C00',
        'gold': 'FFD700', 'tealDecor': '2D2D2D', 'bgCard': '242424', 'textWhite': 'F5F5F5',
        'textBlack': '000000', 'decorPurple': '4A2511'
    }
}

class Config(BaseModel):
    mainTitle1: str
    mainTitle2: str
    pill1: str
    pill2: str
    footer: str

class Option(BaseModel):
    label: str
    text: str

class SlideData(BaseModel):
    badge: str
    tag: str
    qText: str
    options: List[Option]

class GenerateRequest(BaseModel):
    config: Config
    activeSlides: List[SlideData]
    themeId: str
    layoutId: str

class AIParsingRequest(BaseModel):
    rawText: str

# Load the XSLT transformer once
XSLT_PATH = os.path.join(os.path.dirname(__file__), 'MML2OMML.XSL')
try:
    xslt_doc = etree.parse(XSLT_PATH)
    xslt_transformer = etree.XSLT(xslt_doc)
except Exception as e:
    print(f"Warning: Could not load {XSLT_PATH}: {e}")
    xslt_transformer = None

def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def inject_omml_math(paragraph, latex_str: str):
    if not xslt_transformer:
        run = paragraph.add_run()
        run.text = f"${latex_str}$"
        return
        
    try:
        latex_str = latex_str.replace('\\\\', '\\').strip()
        mathml = latex2mathml.converter.convert(latex_str)
        
        # CRITICAL FIX: Robust regex to inject the namespace into the root <math> tag
        # latex2mathml outputs <math display="inline">, not bare <math>,
        # so a simple .replace('<math>', ...) silently fails
        if 'xmlns=' not in mathml:
            mathml = re.sub(r'<math([^>]*)>', r'<math\1 xmlns="http://www.w3.org/1998/Math/MathML">', mathml)
            
        mathml_doc = etree.fromstring(mathml.encode('utf-8'))
        omml_doc = xslt_transformer(mathml_doc)
        omml_root = omml_doc.getroot()
        
        paragraph._p.append(omml_root)
    except Exception as e:
        print(f"OMML Injection failed for '{latex_str}': {e}")
        run = paragraph.add_run()
        run.text = f"${latex_str}$"

def populate_paragraph_with_mixed_content(p, text: str, color: str, font_size: int = 14):
    theme_color = hex_to_rgb(color)
    p.line_spacing = 1.5
    parts = re.split(r'(\$\$[\s\S]*?\$\$|\$[\s\S]*?\$)', text)
    
    for part in parts:
        if not part:
            continue
        if part.startswith('$$') and part.endswith('$$'):
            inject_omml_math(p, part[2:-2])
        elif part.startswith('$') and part.endswith('$'):
            inject_omml_math(p, part[1:-1])
        else:
            run = p.add_run()
            run.text = part
            run.font.name = 'Arial'
            run.font.size = Pt(font_size)
            run.font.color.rgb = theme_color

def add_mixed_content(slide, text: str, start_x: float, start_y: float, width: float, color: str, font_size: int = 14) -> float:
    """Adds a text box, parsing out $...$ or $$...$$ as native OMML math."""
    if not text:
        return start_y
        
    theme_color = hex_to_rgb(color)
    
    # Calculate approximate height
    chars_per_line = int(width / 0.09)
    lines = 0
    for line in text.split('\n'):
        lines += max(1, len(line) // chars_per_line)
    
    # 1.5 line spacing multiplier approximate
    text_height = max(0.2, lines * (font_size * 0.018 * 1.5))
    
    txBox = slide.shapes.add_textbox(Inches(start_x), Inches(start_y), Inches(width), Inches(text_height))
    tf = txBox.text_frame
    tf.margin_top = 0
    tf.margin_left = 0
    tf.margin_bottom = 0
    tf.margin_right = 0
    tf.word_wrap = True
    
    populate_paragraph_with_mixed_content(tf.paragraphs[0], text, color, font_size)

    return start_y + text_height + 0.05

def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color)
    shape.line.fill.background()
    return shape

def add_round_rect(slide, x, y, w, h, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color)
    if border_color:
        shape.line.color.rgb = hex_to_rgb(border_color)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius slightly if possible (python-pptx uses adjustments)
    try:
        shape.adjustments[0] = 0.2
    except:
        pass
    return shape

def add_ellipse(slide, x, y, w, h, color, transparency=0.8):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color)
    shape.line.fill.background()
    # Transparency requires direct oxml manipulation in python-pptx, skipping for simplicity or basic implementation
    return shape

def add_text_basic(slide, text, x, y, w, h, color, font_size, bold=False, align=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.margin_top = 0
    tf.margin_left = 0
    tf.margin_bottom = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return txBox

def render_global_decorations(slide, theme):
    # Ellipses (transparency is hard in raw python-pptx without deep oxml, we'll draw them as light versions)
    # Bottom ribbon
    add_rect(slide, 0, 5.5, 10, 0.125, theme['gold'])
    # Logo
    add_text_basic(slide, "SLIDEGEN PRO", 0.3, 5.2, 3.0, 0.25, theme['cyan'], 10, bold=True)

def build_modern_sidebar_title(slide, config, theme):
    render_global_decorations(slide, theme)
    add_rect(slide, 0, 0, 0.05, 5.625, theme['cyan'])
    add_rect(slide, 0.05, 0, 0.05, 5.625, theme['purple'])
    
    # Title
    add_text_basic(slide, f"{config.mainTitle1} {config.mainTitle2}", 1.2, 1.5, 8, 1.5, theme['cyan'], 54, bold=True)
    
    add_rect(slide, 1.2, 3.0, 6.5, 0.04, theme['gold'])
    
    add_round_rect(slide, 1.2, 3.4, 2.8, 0.5, theme['cyan'])
    add_text_basic(slide, config.pill1, 1.2, 3.4, 2.8, 0.5, theme['textBlack'], 16, bold=True, align=PP_ALIGN.CENTER)
    
    add_round_rect(slide, 4.4, 3.4, 2.8, 0.5, theme['bgColor'], border_color=theme['purple'])
    add_text_basic(slide, config.pill2, 4.4, 3.4, 2.8, 0.5, theme['purple'], 16, bold=True, align=PP_ALIGN.CENTER)
    
    add_text_basic(slide, config.footer, 1.2, 4.2, 8, 0.5, theme['textWhite'], 16)

def build_classic_header_title(slide, config, theme):
    render_global_decorations(slide, theme)
    add_rect(slide, 0, 0, 10, 2, theme['purple'])
    add_rect(slide, 0, 1.9, 10, 0.1, theme['gold'])
    
    add_text_basic(slide, f"{config.mainTitle1} {config.mainTitle2}", 0.5, 0.3, 9, 1.2, theme['textBlack'], 45, bold=True, align=PP_ALIGN.CENTER)
    
    add_round_rect(slide, 2, 2.8, 2.8, 0.5, theme['cyan'])
    add_text_basic(slide, config.pill1, 2, 2.8, 2.8, 0.5, theme['textBlack'], 16, bold=True, align=PP_ALIGN.CENTER)
    
    add_round_rect(slide, 5.2, 2.8, 2.8, 0.5, theme['purple'])
    add_text_basic(slide, config.pill2, 5.2, 2.8, 2.8, 0.5, theme['textBlack'], 16, bold=True, align=PP_ALIGN.CENTER)
    
    add_text_basic(slide, config.footer, 0.5, 4.5, 9, 0.5, theme['textWhite'], 16, align=PP_ALIGN.CENTER)

def build_split_focus_title(slide, config, theme):
    render_global_decorations(slide, theme)
    add_rect(slide, 0, 0, 4.5, 5.625, theme['purple'])
    add_rect(slide, 4.5, 0, 0.05, 5.625, theme['gold'])
    
    add_round_rect(slide, 0.8, 2.0, 2.8, 0.5, theme['cyan'])
    add_text_basic(slide, config.pill1, 0.8, 2.0, 2.8, 0.5, theme['textBlack'], 16, bold=True, align=PP_ALIGN.CENTER)
    
    add_round_rect(slide, 0.8, 2.8, 2.8, 0.5, theme['gold'])
    add_text_basic(slide, config.pill2, 0.8, 2.8, 2.8, 0.5, theme['textBlack'], 16, bold=True, align=PP_ALIGN.CENTER)
    
    add_text_basic(slide, f"{config.mainTitle1}\n{config.mainTitle2}", 5, 1.5, 4.5, 2, theme['textWhite'], 49, bold=True)
    add_text_basic(slide, config.footer, 5, 4.5, 4.5, 0.5, theme['textWhite'], 14)

def render_options_grid(slide, q: SlideData, options_color: str, q_end_y: float, start_x: float, content_width: float):
    if not q.options:
        return
        
    options_y = q_end_y + 0.3
    
    total_chars = sum(len(opt.text) for opt in q.options)
    cols = 4 if (total_chars < 60 and len(q.options) == 4) else 2
    rows = 1 if cols == 4 else ((len(q.options) + 1) // 2)
    
    # Use a table to perfectly distribute options across the width
    table_shape = slide.shapes.add_table(rows, cols, Inches(start_x), Inches(options_y), Inches(content_width), Inches(0.5))
    table = table_shape.table
    
    # Attempt to remove visible borders by clearing fills (if supported by theme)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.background()
            # Remove padding
            cell.margin_left = 0
            cell.margin_right = 0
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            
    for i, opt in enumerate(q.options):
        row_idx = i // cols
        col_idx = i % cols
        cell = table.cell(row_idx, col_idx)
        
        p = cell.text_frame.paragraphs[0]
        populate_paragraph_with_mixed_content(p, f"({opt.label}) {opt.text}", options_color, 14)

def build_modern_sidebar_question(slide, q: SlideData, theme):
    render_global_decorations(slide, theme)
    add_rect(slide, 0, 0, 0.05, 5.625, theme['cyan'])
    add_rect(slide, 0.05, 0, 0.05, 5.625, theme['purple'])
    
    # Badge
    add_round_rect(slide, 0.4, 0.2, 0.6, 0.25, theme['cyan'])
    add_text_basic(slide, q.badge, 0.4, 0.2, 0.6, 0.25, theme['textBlack'], 12, bold=True, align=PP_ALIGN.CENTER)
    
    # Tag bottom right
    add_round_rect(slide, 7.5, 5.1, 2.0, 0.25, theme['bgCard'], border_color=theme['purple'])
    add_text_basic(slide, q.tag.upper(), 7.5, 5.1, 2.0, 0.25, theme['purple'], 10, bold=True, align=PP_ALIGN.CENTER)
    
    start_x, start_y, content_width = 0.4, 0.6, 9.4
    q_end_y = add_mixed_content(slide, q.qText, start_x, start_y, content_width, theme['textWhite'], 14)
    render_options_grid(slide, q, theme['cyan'], q_end_y, start_x, content_width)

def build_classic_header_question(slide, q: SlideData, theme):
    render_global_decorations(slide, theme)
    add_rect(slide, 0, 0, 10, 0.05, theme['purple'])
    add_rect(slide, 0, 0.05, 10, 0.02, theme['gold'])
    
    add_round_rect(slide, 0.4, 0.2, 0.6, 0.25, theme['cyan'])
    add_text_basic(slide, q.badge, 0.4, 0.2, 0.6, 0.25, theme['textBlack'], 12, bold=True, align=PP_ALIGN.CENTER)
    
    add_round_rect(slide, 7.5, 5.1, 2.0, 0.25, theme['bgColor'], border_color=theme['purple'])
    add_text_basic(slide, q.tag.upper(), 7.5, 5.1, 2.0, 0.25, theme['purple'], 10, bold=True, align=PP_ALIGN.CENTER)
    
    start_x, start_y, content_width = 0.4, 0.6, 9.4
    q_end_y = add_mixed_content(slide, q.qText, start_x, start_y, content_width, theme['textWhite'], 14)
    render_options_grid(slide, q, theme['cyan'], q_end_y, start_x, content_width)

def build_split_focus_question(slide, q: SlideData, theme):
    render_global_decorations(slide, theme)
    add_rect(slide, 0, 0, 0.1, 5.625, theme['purple'])
    
    add_round_rect(slide, 0.4, 0.2, 0.6, 0.25, theme['cyan'])
    add_text_basic(slide, q.badge, 0.4, 0.2, 0.6, 0.25, theme['textBlack'], 12, bold=True, align=PP_ALIGN.CENTER)
    
    add_round_rect(slide, 7.5, 5.1, 2.0, 0.25, theme['bgColor'], border_color=theme['purple'])
    add_text_basic(slide, q.tag.upper(), 7.5, 5.1, 2.0, 0.25, theme['purple'], 10, bold=True, align=PP_ALIGN.CENTER)
    
    start_x, start_y, content_width = 0.4, 0.6, 9.4
    q_end_y = add_mixed_content(slide, q.qText, start_x, start_y, content_width, theme['textWhite'], 14)
    render_options_grid(slide, q, theme['gold'], q_end_y, start_x, content_width)


@app.post("/api/generate-pptx")
def generate_pptx(req: GenerateRequest):
    theme = THEMES.get(req.themeId, THEMES['dark-neon'])
    
    prs = Presentation()
    # 16:9 ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    blank_layout = prs.slide_layouts[6] # completely blank
    
    # Title Slide
    title_slide = prs.slides.add_slide(blank_layout)
    # Background color isn't perfectly supported natively on slide master without deep oxml, 
    # so we add a full screen rect for background
    add_rect(title_slide, 0, 0, 10, 5.625, theme['bgColor'])
    
    if req.layoutId == 'classic-header':
        build_classic_header_title(title_slide, req.config, theme)
    elif req.layoutId == 'split-focus':
        build_split_focus_title(title_slide, req.config, theme)
    else:
        build_modern_sidebar_title(title_slide, req.config, theme)
        
    for q in req.activeSlides:
        q_slide = prs.slides.add_slide(blank_layout)
        add_rect(q_slide, 0, 0, 10, 5.625, theme['bgColor'])
        
        if req.layoutId == 'classic-header':
            build_classic_header_question(q_slide, q, theme)
        elif req.layoutId == 'split-focus':
            build_split_focus_question(q_slide, q, theme)
        else:
            build_modern_sidebar_question(q_slide, q, theme)
            
    f = BytesIO()
    prs.save(f)
    f.seek(0)
    
    filename = f"{req.config.mainTitle1}_{req.config.mainTitle2}_Presentation.pptx".replace(" ", "_")
    return StreamingResponse(
        f, 
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

import traceback

@app.post("/api/generate")
def generate_ai(req: AIParsingRequest):
    try:
        if not req.rawText or not req.rawText.strip():
            raise HTTPException(status_code=400, detail="No text provided.")
            
        api_key = os.getenv("VITE_GEMINI_API_KEY") or os.getenv("GEMINI_API_KEY")
        if not api_key:
            raise HTTPException(status_code=500, detail="API Key missing! Please add VITE_GEMINI_API_KEY to your Render environment variables.")
            
        client = genai.Client(api_key=api_key)
        
        prompt = f"""You are an elite educational data formatter for a presentation generator. Parse the raw text into a valid JSON array of question objects.
Each object must have exactly these keys: "badge", "tag", "qText", "options".

- "badge": A short capitalized label (e.g. "Q.1", "EXP.13", "CS.1").
- "tag": A category string based on context headings (e.g. "Practice Question", "Worked Example", "Case Study"). Default to "Practice Question" if ambiguous.
- "qText": The question text as a string.
- "options": The extracted options as an array of objects, where each object has a "label" (e.g., "a", "b", "1") and "text" (the content). If no options exist, return an empty array [].

CRITICAL FORMATTING RULES TO PREVENT SLIDE OVERFLOW:

RULE 1 — SPACE ECONOMY (HIGHEST PRIORITY):
- Presentation slides have limited vertical space. DO NOT use excessive double newlines (\\n\\n).
- Use single newlines (\\n) to separate sub-parts or bullet points within qText to save space.

RULE 2 — MULTI-PART QUESTIONS:
If a question has "Part 1" and "Part 2" (or I/II, or Statement I/II), it is ONE question with TWO sub-parts.
- In "qText": Put both parts separated by a SINGLE newline (\\n).
- In "options": You MUST detect when there are TWO separate sets of (a)(b)(c)(d).
  If it has two separate sets, combine them logically or leave "options" empty and put them in "qText".

RULE 3 — COMPACT OPTIONS (GRID LAYOUT):
- Always extract options into the "options" array.
- ALL options MUST be labeled with lowercase letters: "a", "b", "c", "d". DO NOT use "1", "2", "3", "4" or uppercase letters. Map them if necessary.
- ONLY use this array structure. Do NOT return a string for options.

RULE 4 — MATHEMATICAL NOTATION (CRITICAL):
- Wrap ALL math expressions in single dollar signs: $...$
- Use standard LaTeX commands with SINGLE backslashes inside the dollar signs.
- Fractions: $\\frac{{a}}{{b}}$. Exponents: $x^{{2}}$. Subscripts: $x_{{1}}$. Roots: $\\sqrt{{x}}$.
- Geometry: $\\triangle ABC$, $\\angle BOC = 130^\\circ$
- Example: "The area is $\\frac{{625}}{{36}}$ $cm^{{2}}$" NOT "The area is 625/36 cm^2"
- NEVER output bare caret (^), slash fractions (a/b), or Unicode sub/superscripts as plain text.
- In the JSON string, backslashes must be escaped as \\\\. So $\\frac{{1}}{{2}}$ in JSON becomes "$\\\\frac{{1}}{{2}}$".

RULE 5 — ASSERTION (A) & REASON (R):
- Separate the Assertion text and Reason text with a single \\n in qText.

RULE 6 — CLEAN QUESTION TEXT:
- STRIP out all metadata, filenames, page numbers, or useless headers (e.g., '3. Questions from "35 QA Geometry-3 Q.pdf"'). If it's not a question, option, or valid context, remove it entirely.
- Do NOT include the original question number or label (like "Exp. 13:", "Case 1:", "Q.1") inside `qText`. Strip it completely. The `qText` should start directly with the actual question content.
- IGNORE the original question numbers in the source text. Assign the 'badge' sequentially starting from Q.1, Q.2, Q.3 based on their order in the text.

RULE 7 — JSON SAFETY:
- Escape all newlines as \\n in JSON string values.
- Escape double quotes with backslash.
- Do NOT use actual unescaped line breaks inside JSON string values.

Raw text to format:
{req.rawText}

Respond ONLY with the JSON array. Do not include markdown wrappers like ```json."""

        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt
        )
        
        text = response.text
        text = re.sub(r'```json\s*', '', text, flags=re.IGNORECASE)
        text = re.sub(r'```\s*', '', text).strip()
        
        array_start = text.find('[')
        array_end = text.rfind(']')
        if array_start == -1 or array_end == -1:
            raise HTTPException(status_code=500, detail="AI did not return a valid JSON array.")
            
        json_string = text[array_start:array_end + 1]
        data = json.loads(json_string)
        
        if not isinstance(data, list) or len(data) == 0:
            raise HTTPException(status_code=500, detail="AI returned an empty or invalid array.")
            
        validated = []
        for i, item in enumerate(data):
            options = item.get('options', [])
            if not isinstance(options, list):
                options = []
            
            # Normalize double-escaped LaTeX: \\frac -> \frac, \\sqrt -> \sqrt, etc.
            def normalize_latex(s):
                if not s:
                    return s
                # Fix double-escaped LaTeX commands (e.g. \\angle -> \angle)
                s = re.sub(
                    r'\\\\(frac|sqrt|triangle|angle|circ|pi|theta|alpha|beta|gamma|delta|sum|prod|int|lim|infty|pm|times|div|cdot|leq|geq|neq|approx|equiv|subset|supset|cap|cup|in|notin|forall|exists|nabla|partial|rightarrow|leftarrow|Rightarrow|Leftarrow|text)\b',
                    lambda m: '\\' + m.group(1),
                    s
                )
                # Fix missing spaces after LaTeX geometry commands (KaTeX fails without space)
                # e.g. \triangleABC -> \triangle ABC, \angleA -> \angle A
                s = re.sub(
                    r'\\(triangle|angle|circ)([A-Za-z])',
                    lambda m: '\\' + m.group(1) + ' ' + m.group(2),
                    s
                )
                return s
            
            qText = normalize_latex(item.get('qText', ''))
            cleaned_options = []
            for opt in options:
                if isinstance(opt, dict):
                    cleaned_options.append({
                        "label": opt.get("label", ""),
                        "text": normalize_latex(opt.get("text", ""))
                    })
                    
            validated.append({
                "badge": f"Q.{i + 1}",
                "tag": item.get('tag', 'Practice Question'),
                "qText": qText,
                "options": cleaned_options
            })
            
        return validated
    except Exception as e:
        print(f"AI Formatting Error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e) or "Failed to process text.")

if __name__ == "__main__":
    import uvicorn
    # Allow port to be configured by Render PORT env var, fallback to 8000
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
